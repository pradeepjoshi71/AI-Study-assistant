import os
import logging
from typing import List, Dict, Any

# Optional imports with fallbacks
try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    import docx
except ImportError:
    docx = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

from app.services.ocr import perform_ocr

logger = logging.getLogger(__name__)

def extract_pdf(file_path: str) -> Dict[str, Any]:
    logger.info(f"Extracting PDF: {file_path}")
    pages = []
    metadata = {}
    
    if fitz:
        try:
            with fitz.open(file_path) as doc:
                metadata = doc.metadata or {}
                metadata["page_count"] = doc.page_count
        except Exception as err:
            logger.warning(f"PyMuPDF metadata extraction failed: {err}")

    if pdfplumber:
        try:
            with pdfplumber.open(file_path) as pdf:
                for idx, page in enumerate(pdf.pages):
                    text = page.extract_text() or ""
                    pages.append({
                        "page_number": idx + 1,
                        "text": text
                    })
        except Exception as err:
            logger.error(f"pdfplumber text extraction failed: {err}")
            raise err
    else:
        # Fallback to PyMuPDF if pdfplumber is not available
        if fitz:
            try:
                with fitz.open(file_path) as doc:
                    for idx, page in enumerate(doc):
                        pages.append({
                            "page_number": idx + 1,
                            "text": page.get_text()
                        })
            except Exception as err:
                logger.error(f"PyMuPDF fallback extraction failed: {err}")
                raise err
        else:
            raise ImportError("Neither pdfplumber nor PyMuPDF is installed")

    return {
        "pages": pages,
        "page_count": len(pages),
        "metadata": {
            "title": metadata.get("title", ""),
            "author": metadata.get("author", ""),
            "subject": metadata.get("subject", ""),
            "creator": metadata.get("creator", ""),
            "producer": metadata.get("producer", ""),
        }
    }

def extract_docx(file_path: str) -> Dict[str, Any]:
    logger.info(f"Extracting DOCX: {file_path}")
    if not docx:
        raise ImportError("python-docx is not installed")

    try:
        doc = docx.Document(file_path)
        paragraphs = []
        headings = []

        for p in doc.paragraphs:
            if p.text.strip():
                paragraphs.append(p.text)
                # Check for headings
                if p.style and p.style.name.startswith("Heading"):
                    headings.append(p.text)

        # DOCX doesn't have native pages, treat entire document as a single page
        text = "\n".join(paragraphs)
        
        # Read core document metadata properties
        meta_props = {}
        try:
            core_properties = doc.core_properties
            meta_props = {
                "title": core_properties.title or "",
                "author": core_properties.author or "",
                "created": str(core_properties.created) if core_properties.created else "",
            }
        except Exception as e:
            logger.warning(f"Could not read DOCX properties: {e}")

        return {
            "pages": [{"page_number": 1, "text": text}],
            "page_count": 1,
            "metadata": {
                "headings": headings,
                **meta_props
            }
        }
    except Exception as err:
        logger.error(f"DOCX extraction failed: {err}")
        raise err

def extract_pptx(file_path: str) -> Dict[str, Any]:
    logger.info(f"Extracting PPTX: {file_path}")
    if not Presentation:
        raise ImportError("python-pptx is not installed")

    try:
        prs = Presentation(file_path)
        pages = []

        for idx, slide in enumerate(prs.slides):
            slide_texts = []
            notes_text = ""
            
            # Extract notes if available
            try:
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text
            except Exception:
                pass

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())

            pages.append({
                "page_number": idx + 1,
                "text": "\n".join(slide_texts),
                "notes": notes_text
            })

        return {
            "pages": pages,
            "page_count": len(prs.slides),
            "metadata": {
                "slide_count": len(prs.slides)
            }
        }
    except Exception as err:
        logger.error(f"PPTX extraction failed: {err}")
        raise err

def extract_txt(file_path: str) -> Dict[str, Any]:
    logger.info(f"Extracting TXT: {file_path}")
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()

        return {
            "pages": [{"page_number": 1, "text": text}],
            "page_count": 1,
            "metadata": {}
        }
    except Exception as err:
        logger.error(f"TXT extraction failed: {err}")
        raise err

def extract_image(file_path: str) -> Dict[str, Any]:
    logger.info(f"Extracting Image via OCR: {file_path}")
    try:
        ocr_result = perform_ocr(file_path)
        return {
            "pages": [{"page_number": 1, "text": ocr_result["text"]}],
            "page_count": 1,
            "metadata": {
                "ocr_confidence": ocr_result["confidence"],
                "ocr_details": ocr_result["details"]
            }
        }
    except Exception as err:
        logger.error(f"Image extraction failed: {err}")
        raise err

def extract_document(file_path: str, file_type: str) -> Dict[str, Any]:
    """
    Route extraction requests based on file type.
    """
    ext = file_type.lower().strip(".")
    
    if ext in ("pdf", "application/pdf"):
        return extract_pdf(file_path)
    elif ext in ("docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document"):
        return extract_docx(file_path)
    elif ext in ("pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation"):
        return extract_pptx(file_path)
    elif ext in ("txt", "text/plain"):
        return extract_txt(file_path)
    elif ext in ("png", "jpg", "jpeg", "image/png", "image/jpeg"):
        return extract_image(file_path)
    else:
        # Fallback to extension check
        _, filename_ext = os.path.splitext(file_path)
        filename_ext = filename_ext.lower().strip(".")
        if filename_ext == "pdf":
            return extract_pdf(file_path)
        elif filename_ext == "docx":
            return extract_docx(file_path)
        elif filename_ext == "pptx":
            return extract_pptx(file_path)
        elif filename_ext == "txt":
            return extract_txt(file_path)
        elif filename_ext in ("png", "jpg", "jpeg"):
            return extract_image(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_type} (ext: {filename_ext})")
