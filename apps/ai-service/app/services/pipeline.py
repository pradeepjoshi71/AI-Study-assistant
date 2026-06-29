import os
import re
import tempfile
import logging
from typing import List, Dict, Any, Optional
import urllib.request
import uuid

# Core Python doc/presentation libraries (they should be installed as per requirements.txt)
try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    import docx
except ImportError:
    docx = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

# External libraries for YouTube, Web, OCR, Whisper, Playwright, Readability
# In production, we import them; in dev or if missing, we gracefully degrade with mock output.
try:
    import pytesseract
except ImportError:
    pytesseract = None

try:
    from PIL import Image
except ImportError:
    Image = None

try:
    import yt_dlp
except ImportError:
    yt_dlp = None

try:
    # faster-whisper or standard whisper package
    from faster_whisper import WhisperModel
except ImportError:
    WhisperModel = None

try:
    from playwright.sync_api import sync_playwright
except ImportError:
    sync_playwright = None

try:
    from readability import Document as ReadabilityDoc
except ImportError:
    ReadabilityDoc = None

# Fallback BeautifulSoup for readability if readability package is missing or fails
try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

logger = logging.getLogger(__name__)

class PipelineService:
    """
    PipelineService extracts content from raw sources/keys and returns structured parsed lines.
    Output: List of Dict containing:
      - text: str
      - pageRef: int (1-based index or page number, defaults to 1)
      - sectionTitle: Optional[str]
    """

    def __init__(self, s3_client=None, bucket_name: str = "study-assistant"):
        self.s3_client = s3_client
        self.bucket_name = bucket_name

    def parse_document(
        self,
        doc_id: str,
        storage_key: str,
        mime_type: str,
        local_path: Optional[str] = None
    ) -> List[Dict[str, Any]]:
        """
        Parses a document/URL based on the storage_key and mime_type.
        If local_path is provided, it reads from it; otherwise, it downloads/fetches.
        """
        logger.info(f"Parsing document: id={doc_id}, key={storage_key}, mime={mime_type}")

        # Check if the storage key is a URL
        is_url = storage_key.startswith("http://") or storage_key.startswith("https://")
        is_youtube = is_url and any(domain in storage_key for domain in ["youtube.com", "youtu.be"])

        if is_youtube:
            return self._parse_youtube(storage_key)
        elif is_url:
            return self._parse_web_url(storage_key)

        # File processing: download from storage if local_path is not already given
        cleanup_local = False
        if not local_path:
            temp_dir = tempfile.gettempdir()
            ext = storage_key.split(".")[-1] if "." in storage_key else "txt"
            local_path = os.path.join(temp_dir, f"{doc_id}_{uuid_like()[:8]}.{ext}")
            cleanup_local = True

            if self.s3_client:
                logger.info(f"Downloading key {storage_key} to {local_path} from bucket {self.bucket_name}")
                self.s3_client.download_file(self.bucket_name, storage_key, local_path)
            else:
                # Local development fallback: write mock/read local if exists
                logger.warning(f"No S3 client configured. Simulating mock download to {local_path}")
                with open(local_path, "w", encoding="utf-8") as f:
                    f.write(
                        "AI Study Assistant. This is a local mock file containing simulated raw content "
                        "for testing document ingestion parsing pipeline routers."
                    )

        try:
            # Route to parser by type/extension
            ext = storage_key.split(".")[-1].lower() if "." in storage_key else ""
            mime_lower = mime_type.lower()

            if ext == "pdf" or "pdf" in mime_lower:
                return self._parse_pdf(local_path)
            elif ext == "docx" or "wordprocessing" in mime_lower:
                return self._parse_docx(local_path)
            elif ext == "pptx" or "presentation" in mime_lower:
                return self._parse_pptx(local_path)
            elif ext in ("png", "jpg", "jpeg", "tiff", "bmp") or "image" in mime_lower:
                return self._parse_image_ocr(local_path)
            elif ext in ("txt", "md", "markdown") or "text" in mime_lower:
                return self._parse_plain_text(local_path)
            else:
                # Default fallback
                return self._parse_plain_text(local_path)
        finally:
            if cleanup_local and os.path.exists(local_path):
                try:
                    os.remove(local_path)
                except Exception as e:
                    logger.warning(f"Failed to remove temp file {local_path}: {e}")

    # ── PDF Parser (PyMuPDF / OCR Fallback) ───────────────────────────────────

    def _parse_pdf(self, path: str) -> List[Dict[str, Any]]:
        if not fitz:
            logger.warning("PyMuPDF (fitz) is not installed. Falling back to plain text mock.")
            return [{"text": "PyMuPDF not installed on service.", "pageRef": 1, "sectionTitle": "Fallback"}]

        structured = []
        try:
            with fitz.open(path) as doc:
                for page_idx in range(len(doc)):
                    page = doc[page_idx]
                    page_num = page_idx + 1
                    text = page.get_text().strip()

                    # Detect scan/empty pages and run Tesseract if text is tiny
                    if len(text) < 50 and (pytesseract and Image):
                        logger.info(f"Page {page_num} text length ({len(text)}) is very small. Attempting OCR.")
                        try:
                            pix = page.get_pixmap()
                            temp_img_path = os.path.join(tempfile.gettempdir(), f"page_{page_num}_{uuid_like()[:8]}.png")
                            pix.save(temp_img_path)
                            ocr_text = self._ocr_image(temp_img_path)
                            if os.path.exists(temp_img_path):
                                os.remove(temp_img_path)
                            if ocr_text:
                                text = ocr_text
                        except Exception as ocr_err:
                            logger.error(f"OCR fallback on PDF page {page_num} failed: {ocr_err}")

                    if not text:
                        continue

                    # Attempt basic section detection
                    sections = self._detect_sections_from_text(text)
                    if sections:
                        for sec_title, sec_text in sections:
                            structured.append({
                                "text": sec_text.strip(),
                                "pageRef": page_num,
                                "sectionTitle": sec_title
                            })
                    else:
                        structured.append({
                            "text": text,
                            "pageRef": page_num,
                            "sectionTitle": None
                        })
        except Exception as e:
            logger.error(f"Failed to parse PDF {path}: {e}")
            raise e

        return structured

    # ── Word DOCX Parser ──────────────────────────────────────────────────────

    def _parse_docx(self, path: str) -> List[Dict[str, Any]]:
        if not docx:
            return [{"text": "python-docx not installed.", "pageRef": 1, "sectionTitle": "Fallback"}]

        structured = []
        try:
            doc = docx.Document(path)
            current_section = None
            current_page = 1
            buffer_text = []

            for p in doc.paragraphs:
                text = p.text.strip()
                if not text:
                    continue

                # Check if paragraph is heading
                if p.style and p.style.name.startswith("Heading"):
                    if buffer_text:
                        structured.append({
                            "text": "\n".join(buffer_text),
                            "pageRef": current_page,
                            "sectionTitle": current_section
                        })
                        buffer_text = []
                    current_section = text
                else:
                    # Check for page break characters/runs
                    if "\x0c" in text:
                        if buffer_text:
                            structured.append({
                                "text": "\n".join(buffer_text),
                                "pageRef": current_page,
                                "sectionTitle": current_section
                            })
                            buffer_text = []
                        current_page += 1
                        text = text.replace("\x0c", "")

                    if text:
                        buffer_text.append(text)

            if buffer_text:
                structured.append({
                    "text": "\n".join(buffer_text),
                    "pageRef": current_page,
                    "sectionTitle": current_section
                })
        except Exception as e:
            logger.error(f"Failed to parse DOCX {path}: {e}")
            raise e

        return structured

    # ── PowerPoint PPTX Parser ────────────────────────────────────────────────

    def _parse_pptx(self, path: str) -> List[Dict[str, Any]]:
        if not Presentation:
            return [{"text": "python-pptx not installed.", "pageRef": 1, "sectionTitle": "Fallback"}]

        structured = []
        try:
            prs = Presentation(path)
            for idx, slide in enumerate(prs.slides):
                page_num = idx + 1
                slide_texts = []
                title = None

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        txt = shape.text.strip()
                        # Simple slide title detection
                        if shape.name.lower().startswith("title") or (shape == slide.shapes[0] and len(txt) < 150):
                            title = txt
                        else:
                            slide_texts.append(txt)

                combined = "\n".join(slide_texts)
                structured.append({
                    "text": combined if combined else (title or "Empty Slide"),
                    "pageRef": page_num,
                    "sectionTitle": title
                })
        except Exception as e:
            logger.error(f"Failed to parse PPTX {path}: {e}")
            raise e

        return structured

    # ── Image OCR (Tesseract) Parser ──────────────────────────────────────────

    def _parse_image_ocr(self, path: str) -> List[Dict[str, Any]]:
        text = self._ocr_image(path)
        if not text:
            text = "No text extracted from image OCR."
        return [{
            "text": text,
            "pageRef": 1,
            "sectionTitle": "OCR Image Content"
        }]

    def _ocr_image(self, path: str) -> str:
        if not pytesseract or not Image:
            logger.warning("pytesseract or PIL is not installed. Mocking image OCR.")
            return "Mock OCR: Tesseract dependencies not fully loaded."
        try:
            img = Image.open(path)
            return pytesseract.image_to_string(img).strip()
        except Exception as e:
            logger.error(f"OCR extraction failed: {e}")
            return ""

    # ── Plain Text / Markdown Parser ──────────────────────────────────────────

    def _parse_plain_text(self, path: str) -> List[Dict[str, Any]]:
        structured = []
        try:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read().strip()

            # Handle Markdown split on headers (# Header) or page breaks (--- / \x0c)
            sections = self._detect_sections_from_text(text)
            if sections:
                for idx, (sec_title, sec_text) in enumerate(sections):
                    structured.append({
                        "text": sec_text.strip(),
                        "pageRef": 1,
                        "sectionTitle": sec_title
                    })
            else:
                structured.append({
                    "text": text,
                    "pageRef": 1,
                    "sectionTitle": None
                })
        except Exception as e:
            logger.error(f"Failed to parse TXT/MD {path}: {e}")
            raise e
        return structured

    # ── YouTube URL Parser (yt-dlp + whisper) ─────────────────────────────────

    def _parse_youtube(self, url: str) -> List[Dict[str, Any]]:
        logger.info(f"Extracting YouTube audio & transcribing: {url}")
        if not yt_dlp:
            return [{"text": "yt-dlp not installed on service.", "pageRef": 1, "sectionTitle": "YouTube Transcript Fallback"}]

        temp_dir = tempfile.gettempdir()
        out_tpl = os.path.join(temp_dir, f"yt_audio_{uuid_like()[:8]}.%(ext)s")
        ydl_opts = {
            'format': 'bestaudio/best',
            'outtmpl': out_tpl,
            'postprocessors': [{
                'key': 'FFmpegExtractAudio',
                'preferredcodec': 'mp3',
                'preferredquality': '192',
            }],
            'quiet': True,
        }

        audio_path = None
        try:
            # Download audio
            with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                info = ydl.extract_info(url, download=True)
                title = info.get("title", "YouTube Video")
                # Deduce path
                audio_path = ydl.prepare_filename(info)
                # Correct suffix if needed (FFmpeg changes it to .mp3)
                audio_path = os.path.splitext(audio_path)[0] + ".mp3"

            # Transcribe audio via Whisper
            transcript_text = self._transcribe_audio(audio_path)
            return [{
                "text": transcript_text,
                "pageRef": 1,
                "sectionTitle": f"Video: {title}"
            }]

        except Exception as e:
            logger.error(f"Failed to extract YouTube content: {e}")
            # Graceful degrade with metadata description
            return [{
                "text": f"YouTube video transcription failed. URL: {url}",
                "pageRef": 1,
                "sectionTitle": "Transcription Failure"
            }]
        finally:
            if audio_path and os.path.exists(audio_path):
                try:
                    os.remove(audio_path)
                except Exception:
                    pass

    def _transcribe_audio(self, path: str) -> str:
        if not WhisperModel:
            logger.warning("faster-whisper is not available. Using mock whisper transcript.")
            return "Mock Whisper Transcript: This is a placeholder transcript representing video/audio narrative."

        try:
            # Load model size 'tiny' or 'base' for memory/speed safety
            model = WhisperModel("tiny", device="cpu", compute_type="int8")
            segments, info = model.transcribe(path, beam_size=5)
            text_parts = [seg.text for seg in segments]
            return " ".join(text_parts).strip()
        except Exception as e:
            logger.error(f"Whisper transcription failed: {e}")
            return "Transcription failed due to a processing error."

    # ── Web URL Parser (Playwright + Readability) ─────────────────────────────

    def _parse_web_url(self, url: str) -> List[Dict[str, Any]]:
        logger.info(f"Extracting webpage content: {url}")
        html_content = ""

        if sync_playwright:
            try:
                with sync_playwright() as p:
                    # Run headless browser to fetch javascript-rendered pages
                    browser = p.chromium.launch(headless=True)
                    page = browser.new_page()
                    page.goto(url, wait_until="networkidle", timeout=15000)
                    html_content = page.content()
                    browser.close()
            except Exception as browser_err:
                logger.error(f"Playwright browser fetch failed: {browser_err}. Falling back to urllib.")

        if not html_content:
            try:
                req = urllib.request.Request(
                    url,
                    headers={'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) Chrome/120.0.0.0'}
                )
                with urllib.request.urlopen(req, timeout=10) as response:
                    html_content = response.read().decode('utf-8', errors='ignore')
            except Exception as e:
                logger.error(f"HTTP fallback fetch failed: {e}")
                return [{"text": f"Failed to fetch content from URL: {url}", "pageRef": 1, "sectionTitle": "Web Page Content Check"}]

        # Clean HTML via Readability
        title = "Web Article"
        extracted_text = ""

        if ReadabilityDoc:
            try:
                doc = ReadabilityDoc(html_content)
                title = doc.title()
                clean_html = doc.summary()
                # strip html tags
                if BeautifulSoup:
                    soup = BeautifulSoup(clean_html, "html.parser")
                    extracted_text = soup.get_text(separator="\n").strip()
                else:
                    extracted_text = re.sub('<[^<]+?>', '', clean_html).strip()
            except Exception as read_err:
                logger.error(f"Readability summary cleaning failed: {read_err}")

        # Fallback to BeautifulSoup alone if readability failed
        if not extracted_text and BeautifulSoup:
            try:
                soup = BeautifulSoup(html_content, "html.parser")
                title = soup.title.string if soup.title else "Webpage"
                # Remove script/style tags
                for element in soup(["script", "style", "nav", "footer", "header"]):
                    element.decompose()
                extracted_text = soup.get_text(separator="\n").strip()
            except Exception:
                pass

        if not extracted_text:
            # Crude regex tag stripping fallback
            extracted_text = re.sub('<[^<]+?>', '', html_content).strip()

        # Clean multiple newlines
        extracted_text = re.sub(r'\n+', '\n', extracted_text)

        return [{
            "text": extracted_text,
            "pageRef": 1,
            "sectionTitle": title
        }]

    # ── Formatting/Section Detection Utilities ───────────────────────────────

    def _detect_sections_from_text(self, text: str) -> List[tuple]:
        """
        Splits text by markdown header blocks or potential headers
        """
        # Look for Markdown headers (e.g. ## Section Title)
        header_pattern = r'(^#+\s+.+$)'
        parts = re.split(header_pattern, text, flags=re.MULTILINE)
        
        if len(parts) < 3:
            return []

        sections = []
        current_title = "Introduction"
        
        # Element 0 is the text before the first header, if any
        if parts[0].strip():
            sections.append((current_title, parts[0].strip()))

        for i in range(1, len(parts), 2):
            current_title = parts[i].strip("#").strip()
            content = parts[i+1].strip() if i+1 < len(parts) else ""
            if content:
                sections.append((current_title, content))

        return sections

def uuid_like() -> str:
    import uuid
    return str(uuid.uuid4())
